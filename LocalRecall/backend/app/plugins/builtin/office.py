from pathlib import Path
from app.plugins.base import DataSourcePlugin, ParsedDocument, Chunk


class OfficePlugin(DataSourcePlugin):
    def supported_extensions(self) -> list[str]:
        return [".docx", ".xlsx", ".pptx"]

    def parse(self, file_path: Path) -> ParsedDocument:
        if not self.can_handle(file_path):
            raise ValueError(f"Unsupported: {file_path.suffix}")
        ext = file_path.suffix.lower()
        if ext == ".docx":
            content = self._parse_docx(file_path)
        elif ext == ".xlsx":
            content = self._parse_xlsx(file_path)
        elif ext == ".pptx":
            content = self._parse_pptx(file_path)
        else:
            content = ""
        return ParsedDocument(
            content=content,
            metadata={"file_name": file_path.name},
            source_path=file_path,
            content_type="office",
        )

    def _parse_docx(self, path: Path) -> str:
        from docx import Document
        doc = Document(str(path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _parse_xlsx(self, path: Path) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    def _parse_pptx(self, path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                parts.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
        return "\n\n".join(parts)

    def chunk(self, doc: ParsedDocument) -> list[Chunk]:
        paragraphs = doc.content.split("\n\n")
        chunks = []
        buffer = ""
        idx = 0
        for para in paragraphs:
            if len(buffer) + len(para) > 1500 and buffer:
                chunks.append(Chunk(text=buffer.strip(), metadata=doc.metadata, index=idx))
                idx += 1
                buffer = para
            else:
                buffer += ("\n\n" + para) if buffer else para
        if buffer.strip():
            chunks.append(Chunk(text=buffer.strip(), metadata=doc.metadata, index=idx))
        return chunks
